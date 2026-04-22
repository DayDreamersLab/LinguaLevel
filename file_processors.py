from io import BytesIO
import os

from docx import Document
from pypdf import PdfReader
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

SUPPORTED_EXTENSIONS = {".txt", ".pdf", ".docx", ".pptx"}


def extract_text(file_stream, extension: str) -> str:
    data = file_stream.read()
    stream = BytesIO(data)

    if extension == ".txt":
        return data.decode("utf-8", errors="ignore")

    if extension == ".pdf":
        reader = PdfReader(stream)
        parts = []
        for page in reader.pages:
            parts.append(page.extract_text() or "")
        return "\n\n".join(parts)

    if extension == ".docx":
        document = Document(stream)
        return "\n".join(p.text for p in document.paragraphs)

    if extension == ".pptx":
        presentation = Presentation(stream)
        collected = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    collected.append(shape.text)
        return "\n".join(collected)

    raise ValueError(f"Unsupported extension {extension}")


def rebuild_file(
    original_filename: str,
    extension: str,
    simplified_text: str,
    target_level: str,
    vocab_changes: list[dict[str, str]],
) -> tuple[BytesIO, str]:
    output_name = f"{os.path.splitext(original_filename)[0]}_simplified_{target_level}{extension}"

    if extension == ".txt":
        body = compose_output_text(simplified_text, vocab_changes)
        return BytesIO(body.encode("utf-8")), output_name

    if extension == ".docx":
        doc = Document()
        doc.add_heading(f"Simplified to CEFR {target_level}", level=1)
        doc.add_paragraph(simplified_text)
        append_vocab_to_doc(doc, vocab_changes)
        stream = BytesIO()
        doc.save(stream)
        stream.seek(0)
        return stream, output_name

    if extension == ".pptx":
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Simplified to CEFR {target_level}"
        slide.placeholders[1].text = simplified_text[:3000]

        vocab_slide = prs.slides.add_slide(prs.slide_layouts[1])
        vocab_slide.shapes.title.text = "Vocabulary Simplified"
        vocab_text = "\n".join(
            f"• {v['original']} → {v['replacement']}: {v['meaning']}" for v in vocab_changes
        )
        vocab_slide.placeholders[1].text = vocab_text[:3000] if vocab_text else "No vocabulary was replaced."

        stream = BytesIO()
        prs.save(stream)
        stream.seek(0)
        return stream, output_name

    if extension == ".pdf":
        stream = BytesIO()
        pdf = canvas.Canvas(stream, pagesize=letter)
        width, height = letter
        y = height - 60

        pdf.setFont("Helvetica-Bold", 14)
        pdf.drawString(50, y, f"Simplified to CEFR {target_level}")
        y -= 30

        pdf.setFont("Helvetica", 10)
        for line in wrap_lines(simplified_text, 100):
            if y < 60:
                pdf.showPage()
                y = height - 60
                pdf.setFont("Helvetica", 10)
            pdf.drawString(50, y, line)
            y -= 14

        y -= 10
        pdf.setFont("Helvetica-Bold", 12)
        if y < 60:
            pdf.showPage()
            y = height - 60
        pdf.drawString(50, y, "Vocabulary Simplified")
        y -= 20

        pdf.setFont("Helvetica", 10)
        vocab_lines = [
            f"- {v['original']} -> {v['replacement']}: {v['meaning']}" for v in vocab_changes
        ] or ["No vocabulary was replaced."]

        for line in vocab_lines:
            for wrapped in wrap_lines(line, 100):
                if y < 60:
                    pdf.showPage()
                    y = height - 60
                    pdf.setFont("Helvetica", 10)
                pdf.drawString(50, y, wrapped)
                y -= 14

        pdf.save()
        stream.seek(0)
        return stream, output_name

    raise ValueError(f"Unsupported extension {extension}")


def wrap_lines(text: str, max_chars: int) -> list[str]:
    words = text.replace("\r", "").split()
    if not words:
        return [""]

    lines = []
    current = []
    length = 0
    for word in words:
        next_len = length + len(word) + (1 if current else 0)
        if next_len > max_chars:
            lines.append(" ".join(current))
            current = [word]
            length = len(word)
        else:
            current.append(word)
            length = next_len
    if current:
        lines.append(" ".join(current))
    return lines


def compose_output_text(simplified_text: str, vocab_changes: list[dict[str, str]]) -> str:
    section = ["Simplified Material", "=" * 20, simplified_text, "", "Vocabulary Simplified", "=" * 20]
    if not vocab_changes:
        section.append("No vocabulary was replaced.")
    else:
        for item in vocab_changes:
            section.append(f"- {item['original']} -> {item['replacement']}: {item['meaning']}")
    return "\n".join(section)


def append_vocab_to_doc(doc: Document, vocab_changes: list[dict[str, str]]) -> None:
    doc.add_heading("Vocabulary Simplified", level=2)
    if not vocab_changes:
        doc.add_paragraph("No vocabulary was replaced.")
        return
    for item in vocab_changes:
        doc.add_paragraph(
            f"{item['original']} → {item['replacement']}: {item['meaning']}",
            style="List Bullet",
        )
